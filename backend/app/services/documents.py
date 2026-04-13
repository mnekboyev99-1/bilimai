from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path

from fastapi import UploadFile
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "txt", "md"}
STOPWORDS = {
    "uz": {"va", "bu", "uchun", "bilan", "ham", "bir", "ikki", "kabi", "yoki", "emas", "shu", "unda", "bor", "qilish", "bo‘ladi", "haqida", "kurs", "dars", "modul"},
    "ru": {"и", "в", "на", "с", "по", "для", "как", "это", "что", "или", "не", "из", "к", "мы", "вы", "они", "курс", "урок", "модуль"},
    "en": {"the", "and", "for", "with", "that", "this", "from", "into", "your", "you", "are", "can", "will", "how", "what", "lesson", "module", "course"},
}


def clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_keywords(text: str, lang: str = "uz", top_n: int = 20) -> list[str]:
    tokens = re.findall(r"[A-Za-zА-Яа-яЁёЎўҚқҒғҲҳʻ’`'-]{3,}", text.lower())
    stop = STOPWORDS.get(lang, STOPWORDS["uz"])
    freq: dict[str, int] = {}
    for token in tokens:
        if token in stop:
            continue
        freq[token] = freq.get(token, 0) + 1
    return [word for word, _ in sorted(freq.items(), key=lambda kv: kv[1], reverse=True)[:top_n]]


def chunk_text(text: str, chunk_size: int = 1800, overlap: int = 240) -> list[str]:
    paragraphs = [p.strip() for p in re.split(r"\n{1,2}", text) if p.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= chunk_size:
            current = f"{current}\n\n{paragraph}".strip()
            continue
        if current:
            chunks.append(current)
            tail = current[-overlap:] if len(current) > overlap else current
            current = f"{tail}\n\n{paragraph}".strip()
        else:
            chunks.append(paragraph[:chunk_size])
            current = paragraph[chunk_size - overlap :] if len(paragraph) > chunk_size else ""
    if current:
        chunks.append(current)
    return chunks


async def parse_upload(file: UploadFile) -> tuple[str, dict]:
    extension = Path(file.filename or "").suffix.lower().replace(".", "")
    raw = await file.read()
    if extension == "pdf":
        reader = PdfReader(BytesIO(raw))
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"[PDF {i}-bet]\n{text}")
        text = "\n\n".join(pages)
        meta = {"pages": len(reader.pages)}
    elif extension == "docx":
        doc = Document(BytesIO(raw))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        text = "\n".join(paragraphs)
        meta = {"paragraphs": len(paragraphs)}
    elif extension == "pptx":
        prs = Presentation(BytesIO(raw))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text.strip())
            if parts:
                slides.append(f"[PPTX {i}-slayd]\n" + "\n".join(parts))
        text = "\n\n".join(slides)
        meta = {"slides": len(prs.slides)}
    else:
        text = raw.decode("utf-8", errors="ignore")
        meta = {"chars": len(text)}
    return clean_text(text), meta


def build_source_summary(chunks: list[str], language: str) -> dict:
    joined = "\n".join(chunks[:8])
    keywords = extract_keywords(joined, language, top_n=18)
    preview = joined[:1200]
    return {"keywords": keywords, "preview": preview, "chunk_count": len(chunks)}
